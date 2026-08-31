import io
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Pt

class PPTXProcessor:
    @staticmethod
    def extract_slides_text(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Trích xuất nội dung văn bản của từng slide kèm số thứ tự slide.
        """
        prs = Presentation(io.BytesIO(file_bytes))
        slides_data = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            # Đọc nội dung trên các khung văn bản (text boxes / shapes)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_texts.append(text)

            content = "\n".join(slide_texts).strip()
            if content:
                slides_data.append({
                    "slide_number": idx,
                    "content": content
                })
        
        return slides_data

    @staticmethod
    def format_doc_text_for_ai(slides_data: List[Dict[str, Any]]) -> str:
        """Định dạng dữ liệu các slide thành văn bản phân tách rõ ràng để AI phân tích."""
        formatted_parts = []
        for item in slides_data:
            formatted_parts.append(
                f"=== SLIDE {item['slide_number']} ===\n{item['content']}\n"
            )
        return "\n".join(formatted_parts)

    @staticmethod
    def integrate_into_notes(file_bytes: bytes, ai_data: dict) -> io.BytesIO:
        """
        Chèn gợi ý Năng lực số / Năng lực AI vào phần Speaker Notes của từng slide.
        """
        prs = Presentation(io.BytesIO(file_bytes))
        sua_doi_list = ai_data.get('sua_doi', [])
        
        # Gom nhóm các nội dung tích hợp theo số thứ tự slide
        notes_by_slide = {}
        for item in sua_doi_list:
            slide_num = item.get('slide_number')
            content = item.get('insert_content', '').strip()
            loai = item.get('loai', 'Năng lực số')
            
            if slide_num is not None and content:
                tag = "[NĂNG LỰC AI]" if loai == "Năng lực AI" else "[NĂNG LỰC SỐ]"
                note_entry = f"{tag}: {content}"
                notes_by_slide.setdefault(int(slide_num), []).append(note_entry)

        # Duyệt qua các slide và chèn vào Notes
        for idx, slide in enumerate(prs.slides, start=1):
            if idx in notes_by_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                
                existing_text = text_frame.text.strip()
                integration_text = "\n\n".join(notes_by_slide[idx])
                
                header_sep = "📌 [TÍCH HỢP NĂNG LỰC SỐ & AI]:\n"
                
                if existing_text:
                    text_frame.text = f"{existing_text}\n\n{header_sep}{integration_text}"
                else:
                    text_frame.text = f"{header_sep}{integration_text}"

        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return output_stream
